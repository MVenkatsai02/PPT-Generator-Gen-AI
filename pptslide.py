import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches

st.set_page_config(page_title="AI-Based Custom Slide Deck Generator", layout="wide")

# Sidebar for API Key Upload
st.sidebar.title("🔑 Upload API Key")
st.sidebar.markdown("""
- [Get Google Gemini API Key](https://aistudio.google.com/app/apikey)  
""")

# API Key Input
gemini_api_key = st.sidebar.text_input("Google Gemini API Key", type="password")

# Ensure API key is provided
if not gemini_api_key:
    st.sidebar.warning("Please enter your API key to proceed.")
    st.stop()

# Initialize Gemini API
genai.configure(api_key=gemini_api_key)

# Streamlit App Main Interface
st.title("📊 AI-Based Custom Slide Deck Generator")
st.subheader("Create professional slide decks instantly with AI!")

# User Inputs
presentation_title = st.text_input("Enter Presentation Title:", "AI in Healthcare")
num_slides = st.slider("Number of Slides:", min_value=3, max_value=15, value=5)
theme = st.selectbox("Select Slide Theme:", ["Professional", "Minimalist", "Creative"])
include_images = st.checkbox("Include Suggested Images")

# Function to generate slide content
def generate_slide_content(title, num_slides, theme, include_images):
    prompt = f"""
    Generate a professional PowerPoint presentation with the title: "{title}".
    Create {num_slides} slides with key points, insights, and structured content.
    Use a "{theme}" theme.
    {"Suggest relevant images for each slide." if include_images else ""}
    Provide slide titles and bullet points for each slide.
    """
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(prompt)
    return response.text if response else "Sorry, I couldn't generate the slide content."

# Function to create PowerPoint slides
def create_presentation(title, slides_content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide_layout = prs.slide_layouts[1]

    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "AI-Generated Presentation"

    # Content Slides
    slides = slides_content.split("\n\n")  # Splitting slides by double line breaks
    for slide_text in slides:
        slide_lines = slide_text.split("\n")
        if len(slide_lines) < 2:
            continue  # Skip empty slides
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_lines[0]  # Slide Title
        content = "\n".join(slide_lines[1:])  # Slide Content
        slide.placeholders[1].text = content

    return prs

# Generate Slide Deck
if st.button("Generate Slide Deck"):
    with st.spinner("Generating slides..."):
        slides_content = generate_slide_content(presentation_title, num_slides, theme, include_images)
        ppt = create_presentation(presentation_title, slides_content)

    # Save and Offer Download
    ppt_filename = f"{presentation_title.replace(' ', '_')}.pptx"
    ppt.save(ppt_filename)
    
    with open(ppt_filename, "rb") as file:
        st.download_button(
            label="📥 Download Slide Deck",
            data=file,
            file_name=ppt_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

# Run the app using:
# streamlit run slide_deck_generator.py
